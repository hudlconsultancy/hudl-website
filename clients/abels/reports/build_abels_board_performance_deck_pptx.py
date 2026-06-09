#!/usr/bin/env python3
"""
Abels — Board Performance Deck builder (PPTX)
Opening / landscape slides that set the tone before the performance data.

Usage:  python3 build_abels_board_performance_deck_pptx.py
Output: 2026-06-abels-board-performance-deck.pptx (same folder)

HUDL brand: dark theme, four-colour stripe, Barlow Condensed look via Arial Narrow.
This deck is board-facing (Abels/AGM). Slides added incrementally.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand palette ─────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x0C, 0x0C, 0x0C)
CARD   = RGBColor(0x14, 0x14, 0x16)
ROW1   = RGBColor(0x0F, 0x0F, 0x12)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SOFT   = RGBColor(0xE6, 0xE6, 0xE6)
YELLOW = RGBColor(0xFB, 0xBF, 0x46)
BLUE   = RGBColor(0x45, 0xAE, 0xFF)
CYAN   = RGBColor(0x1E, 0xDF, 0xD4)
RED    = RGBColor(0xFF, 0x51, 0x5E)
MUTED  = RGBColor(0x9A, 0x9A, 0x9A)

FONT_DISPLAY = "Arial Narrow"
FONT_BODY    = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = BLACK
    return s

def stripe(s):
    bw = SW / 4
    for i, col in enumerate([YELLOW, BLUE, RED, CYAN]):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(i*bw)), 0, Emu(int(bw)), Pt(6))
        r.fill.solid(); r.fill.fore_color.rgb = col
        r.line.fill.background(); r.shadow.inherit = False

def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None: r.line.fill.background()
    else: r.line.color.rgb = line; r.line.width = Pt(line_w)
    r.shadow.inherit = False
    return r

def text(s, txt, l, t, w, h, size=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         font=FONT_BODY, anchor=MSO_ANCHOR.TOP, line_spacing=None, wrap=True):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = font
    return tb

def footer(s):
    text(s, "Abels Moving Services  ·  Board Performance Review  ·  June 2026",
         Inches(0.7), Inches(7.04), Inches(12), Inches(0.32),
         size=9, color=MUTED, font=FONT_BODY)


# ══ SLIDE — Why Digital Marketing Matters ═════════════════════════════════════
s = slide(); stripe(s)

# kicker + left accent bar
rect(s, Inches(0.7), Inches(0.95), Inches(0.14), Inches(0.95), fill=YELLOW)
text(s, "SETTING THE LANDSCAPE", Inches(1.0), Inches(0.92), Inches(10), Inches(0.35),
     size=13, color=YELLOW, bold=True, font=FONT_BODY)
text(s, "Why Digital Marketing Matters", Inches(0.98), Inches(1.22), Inches(11.6), Inches(0.95),
     size=40, color=WHITE, bold=True, font=FONT_DISPLAY)

# body paragraph on a subtle panel
rect(s, Inches(0.7), Inches(2.55), Inches(11.95), Inches(2.7), fill=ROW1)
rect(s, Inches(0.7), Inches(2.55), Inches(0.1), Inches(2.7), fill=BLUE)
para = ("Moving home or abroad is a major, heavily-researched decision that now almost always "
        "starts online — so digital is where a premium brand like Abels is found, judged and "
        "shortlisted. Every pound spent there is fully accountable, tied directly to enquiries, "
        "quotes and booked jobs, which is something traditional advertising can never offer. And "
        "it is where the battle for high-value customers is won: it builds the top of the funnel "
        "the whole business depends on.")
text(s, para, Inches(1.05), Inches(2.75), Inches(11.3), Inches(2.3),
     size=18, color=SOFT, font=FONT_BODY, line_spacing=1.25, anchor=MSO_ANCHOR.MIDDLE)

# funnel pull-quote band
rect(s, Inches(0.7), Inches(5.55), Inches(11.95), Inches(0.95), fill=BLACK, line=YELLOW, line_w=1.25)
text(s, "Visibility  →  Enquiries  →  Quotes  →  Booked jobs  →  Revenue",
     Inches(0.9), Inches(5.55), Inches(11.6), Inches(0.95),
     size=19, color=YELLOW, bold=True, font=FONT_DISPLAY,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

footer(s)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "2026-06-abels-board-performance-deck.pptx")
prs.save(out)
print("Saved:", out, "—", len(prs.slides._sldIdLst), "slide(s)")
