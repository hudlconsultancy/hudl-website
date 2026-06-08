#!/usr/bin/env python3
"""
Abels Moving Services — Board Report (What / Why / How) builder (PPTX)
Generates a real PowerPoint that opens cleanly in Google Slides.

Usage:  python3 build_abels_what_why_how_pptx.py
Output: 2026-06-08-abels-board-report-what-why-how.pptx (same folder)

Source of all figures: 2026-06-08-abels-board-report-what-why-how.md
HUDL brand: dark theme, four-colour stripe, Barlow Condensed look via Arial Narrow.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand palette ─────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x0C, 0x0C, 0x0C)
CARD   = RGBColor(0x14, 0x14, 0x16)
ROW1   = RGBColor(0x0F, 0x0F, 0x12)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFB, 0xBF, 0x46)
BLUE   = RGBColor(0x45, 0xAE, 0xFF)
CYAN   = RGBColor(0x1E, 0xDF, 0xD4)
RED    = RGBColor(0xFF, 0x51, 0x5E)
GREEN  = RGBColor(0x22, 0xC5, 0x5E)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
MUTED  = RGBColor(0x9A, 0x9A, 0x9A)
TBLHDR = RGBColor(0x0A, 0x1A, 0x2A)

FONT_DISPLAY = "Arial Narrow"
FONT_BODY    = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = BLACK
    return s

def stripe(s):
    bw = SW / 4
    for i, col in enumerate([YELLOW, BLUE, RED, CYAN]):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(i*bw)), 0, Emu(int(bw)), Pt(6))
        r.fill.solid(); r.fill.fore_color.rgb = col
        r.line.fill.background(); r.shadow.inherit = False

def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None: r.line.fill.background()
    else: r.line.color.rgb = line; r.line.width = Pt(line_w)
    r.shadow.inherit = False
    return r

def text(s, txt, l, t, w, h, size=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         font=FONT_BODY, anchor=MSO_ANCHOR.TOP, wrap=True, space=None):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space is not None: p.space_after = Pt(space)
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = font
    return tb

def footer(s, page):
    text(s, "Abels Moving Services  ·  Jan–May 2026  ·  Internal — not for external distribution",
         Inches(0.5), Inches(7.06), Inches(10), Inches(0.3), size=8, color=MUTED, font=FONT_BODY)
    text(s, str(page), Inches(12.6), Inches(7.06), Inches(0.5), Inches(0.3),
         size=8, color=MUTED, align=PP_ALIGN.RIGHT)

def heading(s, kicker, title, accent=YELLOW):
    rect(s, Inches(0.5), Inches(0.42), Inches(0.12), Inches(0.62), fill=accent)
    text(s, kicker, Inches(0.78), Inches(0.40), Inches(11), Inches(0.3),
         size=12, color=accent, bold=True, font=FONT_BODY)
    text(s, title, Inches(0.78), Inches(0.66), Inches(11.8), Inches(0.6),
         size=27, color=WHITE, bold=True, font=FONT_DISPLAY)

def add_table(s, l, t, w, h, headers, rows, col_w, row_fills=None,
              header_size=12, cell_size=11, cell_colors=None):
    nrows, ncols = len(rows) + 1, len(headers)
    gf = s.shapes.add_table(nrows, ncols, l, t, w, h); tbl = gf.table
    tbl.first_row = False; tbl.horz_banding = False
    total = sum(col_w)
    for c, cw in enumerate(col_w):
        tbl.columns[c].width = Emu(int(w * cw / total))
    for c, htxt in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = TBLHDR
        cell.margin_left = Pt(7); cell.margin_right = Pt(7)
        cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.LEFT
        run = p.add_run(); run.text = htxt
        run.font.size = Pt(header_size); run.font.bold = True
        run.font.color.rgb = BLUE; run.font.name = FONT_BODY
    for r, row in enumerate(rows):
        fill = ROW1 if r % 2 == 0 else BLACK
        if row_fills and r in row_fills: fill = row_fills[r]
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
            cell.margin_left = Pt(7); cell.margin_right = Pt(7)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = val
            run.font.size = Pt(cell_size)
            col = WHITE
            if cell_colors and (r, c) in cell_colors: col = cell_colors[(r, c)]
            run.font.color.rgb = col; run.font.name = FONT_BODY
            run.font.bold = (c == 0)
    return tbl

def wwh_slide(page, kicker, title, accent, numbers, what, why, how, nxt):
    """Channel slide: headline numbers band + three What/Why/How panels."""
    s = slide(); stripe(s); heading(s, kicker, title, accent)
    # numbers band
    rect(s, Inches(0.5), Inches(1.45), Inches(12.33), Inches(0.7), fill=CARD)
    rect(s, Inches(0.5), Inches(1.45), Inches(0.08), Inches(0.7), fill=accent)
    text(s, numbers, Inches(0.78), Inches(1.45), Inches(12), Inches(0.7),
         size=14, color=WHITE, bold=True, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)
    # three panels
    panels = [("WHAT WE DID", what, accent), ("WHY WE DID IT", why, BLUE), ("HOW IT'S GOING", how, CYAN)]
    pw = Inches(3.94); gap = Inches(0.255); x0 = Inches(0.5); top = Inches(2.45); ph = Inches(3.5)
    for i, (label, body, col) in enumerate(panels):
        x = Emu(int(x0) + i * (int(pw) + int(gap)))
        rect(s, x, top, pw, ph, fill=ROW1)
        rect(s, x, top, pw, Inches(0.5), fill=col)
        text(s, label, x, top, pw, Inches(0.5), size=12.5, color=BLACK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_BODY)
        text(s, body, Emu(int(x)+Pt(6)), Emu(int(top)+int(Inches(0.62))), Emu(int(pw)-Pt(12)),
             Inches(2.8), size=12.5, color=WHITE, font=FONT_BODY, space=6)
    # next strip
    rect(s, Inches(0.5), Inches(6.15), Inches(12.33), Inches(0.62), fill=BLACK, line=accent, line_w=1.25)
    text(s, "NEXT", Inches(0.7), Inches(6.15), Inches(1.1), Inches(0.62),
         size=12, color=accent, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=FONT_BODY)
    text(s, nxt, Inches(1.7), Inches(6.15), Inches(11), Inches(0.62),
         size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font=FONT_BODY)
    footer(s, page)
    return s


# ══ SLIDE 1 — Title ═══════════════════════════════════════════════════════════
s = slide(); stripe(s)
rect(s, Inches(0.9), Inches(2.5), Inches(0.16), Inches(2.0), fill=YELLOW)
text(s, "ABELS MOVING SERVICES", Inches(1.25), Inches(2.45), Inches(11), Inches(0.6),
     size=20, color=MUTED, bold=True, font=FONT_BODY)
text(s, "Performance Report", Inches(1.2), Inches(2.95), Inches(11.5), Inches(1.0),
     size=50, color=WHITE, bold=True, font=FONT_DISPLAY)
text(s, "The What, Why & How  ·  January–May 2026", Inches(1.25), Inches(4.05), Inches(11.5), Inches(0.6),
     size=22, color=YELLOW, bold=True, font=FONT_DISPLAY)
text(s, "Paid Search  ·  Website  ·  SEO", Inches(1.25), Inches(4.7), Inches(11), Inches(0.4),
     size=14, color=MUTED, font=FONT_BODY)
text(s, "Prepared for the board  ·  HUDL  ·  Internal — not for external distribution",
     Inches(1.25), Inches(6.4), Inches(11), Inches(0.4), size=11, color=MUTED, font=FONT_BODY)

# ══ SLIDE 2 — At a glance ═════════════════════════════════════════════════════
s = slide(); stripe(s); heading(s, "AT A GLANCE", "Three channels, one picture", YELLOW)
add_table(s, Inches(0.5), Inches(1.6), Inches(12.33), Inches(2.1),
          ["Channel", "The number that matters", "Trend"],
          [["Paid Search (PPC)", "£23,375 spend  →  136 enquiries  at  £172 each", "Settled to ~£160/lead in Apr–May"],
           ["Website (GA4)", "15,341 visits (Jan–May)", "Eased from Jan; engagement up"],
           ["SEO (organic)", "4,368 clicks from 1.24m search appearances", "Stable but flat; brand-dependent"]],
          col_w=[2.6, 6.2, 3.4],
          cell_colors={(0,0):YELLOW,(1,0):BLUE,(2,0):CYAN},
          header_size=12, cell_size=12)
# story band
rect(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(2.4), fill=ROW1)
rect(s, Inches(0.5), Inches(4.1), Inches(0.1), Inches(2.4), fill=YELLOW)
text(s, "THE ONE-LINE STORY", Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.4),
     size=12, color=YELLOW, bold=True, font=FONT_BODY)
text(s, "Paid Search was rebuilt in spring and is now running efficiently — about £160 per genuine enquiry.\n"
        "The website is quieter than January, but the people on it are more engaged.\n"
        "Organic search has huge visibility but isn't yet converting it — and the opportunity is clearly named in this pack.",
     Inches(0.8), Inches(4.75), Inches(11.6), Inches(1.6), size=15, color=WHITE, font=FONT_BODY, space=8)
footer(s, 2)

# ══ SLIDE 3 — Paid Search ═════════════════════════════════════════════════════
wwh_slide(
    3, "1 · PAID SEARCH (PPC)", "Rebuilt, and now efficient", YELLOW,
    "£23,375 spent     ·     136 enquiries     ·     £172 average cost per enquiry     ·     1.9× return on spend (Apr–May)",
    "Rebuilt the Google Ads account in March around four clear campaigns "
    "(Europe, London, Domestic, International), switched off Performance Max, "
    "and brought management in-house during the JDR review.",
    "Performance Max had spent £1,010 over Jan–Feb for ZERO enquiries. The old "
    "setup also counted soft actions (page views, phone clicks) as conversions, "
    "which flattered the numbers. We wanted real enquiries, properly measured.",
    "April & May are the true picture: ~£6k a month, ~37 enquiries, ~£160 each. "
    "Europe and London are strongest (£110–138) with room to grow; Domestic is "
    "dearest but improving (£436 → £199 in May).",
    "Move more budget into Europe & London; keep tightening Domestic.")

# ══ SLIDE 4 — Website ═════════════════════════════════════════════════════════
wwh_slide(
    4, "2 · WEBSITE", "Quieter, but more engaged", BLUE,
    "15,341 visits     ·     12,061 new visitors     ·     ~44s average engagement     ·     best engagement of the period in May",
    "Tracked website traffic and engagement month by month, alongside the paid "
    "and organic activity, so spend connects to on-site behaviour.",
    "The website is where every enquiry is made — its traffic and engagement are "
    "the bridge between marketing spend and actual leads.",
    "Traffic eased from a January peak (~3,900 visits) to a steady ~2,600–2,850 a "
    "month — expected, as we stopped buying low-intent PMAX traffic. Crucially, "
    "engagement HELD and hit its best level in May.",
    "Confirm which on-site actions we count as enquiries; supply a channel-split export (paid vs organic vs direct).")

# ══ SLIDE 5 — SEO ═════════════════════════════════════════════════════════════
wwh_slide(
    5, "3 · SEO (ORGANIC SEARCH)", "Huge visibility, not yet converted", CYAN,
    "4,368 clicks     ·     1,242,569 appearances in Google     ·     0.35% click rate     ·     average position ~25 (page 2–3)",
    "Abels appears in Google an enormous amount (1.24m times in five months) but "
    "earns relatively few clicks — because most pages rank on page 2–3, not page 1.",
    "The demand is already there. The big commercial pages (international removals, "
    "London removals, London storage) hold 27% of all our search visibility but rank "
    "page 3–5. Today 42% of organic clicks are just the homepage — people already "
    "searching 'Abels' by name.",
    "Focus SEO on lifting those high-visibility commercial pages from page 3–5 onto "
    "page 1–2. The demand is proven; the rankings are the gap. Our guides and the "
    "Royal Warrant page already rank well and support credibility.",
    "A focused SEO push on the London / international / storage pages — the biggest single opportunity.")

# ══ SLIDE 6 — Initiatives executed ════════════════════════════════════════════
s = slide(); stripe(s); heading(s, "INITIATIVES EXECUTED", "What we did, why, and where it stands", YELLOW)
add_table(s, Inches(0.5), Inches(1.6), Inches(12.33), Inches(4.6),
          ["Initiative", "What", "Why", "Status"],
          [["Paused Performance Max", "Switched PMAX off in March", "£1,010 spent, 0 enquiries", "Done"],
           ["Rebuilt the Search account", "Four clear campaigns from March", "Control + real enquiry tracking", "Done — CPA ~£160"],
           ["Fixed conversion tracking", "Count genuine enquiries, not soft clicks", "Honest, decision-grade numbers", "Done"],
           ["Budget reallocation", "Move spend to Europe & London", "They deliver the cheapest enquiries", "In progress"],
           ["Domestic clean-up", "Tighten the most expensive campaign", "CPA was £200–436", "In progress — May £199"],
           ["SEO opportunity identified", "Lift commercial pages off page 3–5", "27% of visibility, almost no clicks", "Next priority"]],
          col_w=[2.9, 3.7, 3.5, 2.2],
          cell_colors={(0,3):GREEN,(1,3):GREEN,(2,3):GREEN,(3,3):AMBER,(4,3):AMBER,(5,3):CYAN},
          header_size=12, cell_size=11)
footer(s, 6)

# ══ SLIDE 7 — Board takeaways ═════════════════════════════════════════════════
s = slide(); stripe(s); heading(s, "FOR THE BOARD", "Three things to note", YELLOW)
items = [
    ("1", "Paid Search is now efficient and honestly measured", YELLOW,
     "~£160 per genuine enquiry, with the best campaigns still able to scale."),
    ("2", "The clearest growth opportunity is SEO", CYAN,
     "The demand exists — we simply need to rank better on the commercial pages."),
    ("3", "One quick win to re-enable: the Moveware → Google Ads link", BLUE,
     "Active in 2025, off in 2026. Switching it back on lets Google optimise toward real booked revenue, not just form-fills."),
]
y = Inches(1.75)
for num, head, col, body in items:
    rect(s, Inches(0.5), y, Inches(12.33), Inches(1.45), fill=ROW1)
    rect(s, Inches(0.5), y, Inches(0.12), Inches(1.45), fill=col)
    text(s, num, Inches(0.75), y, Inches(1.0), Inches(1.45), size=40, color=col, bold=True,
         font=FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, head, Inches(1.85), Emu(int(y)+int(Inches(0.22))), Inches(10.7), Inches(0.55),
         size=17, color=WHITE, bold=True, font=FONT_DISPLAY)
    text(s, body, Inches(1.85), Emu(int(y)+int(Inches(0.78))), Inches(10.7), Inches(0.55),
         size=12.5, color=MUTED, font=FONT_BODY)
    y = Emu(int(y) + int(Inches(1.62)))
footer(s, 7)


out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "2026-06-08-abels-board-report-what-why-how.pptx")
prs.save(out)
print("Saved:", out, "—", len(prs.slides._sldIdLst), "slides")
