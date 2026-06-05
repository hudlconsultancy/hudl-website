#!/usr/bin/env python3
"""
Abels Moving Services — Board Presentation builder (PPTX)
Generates a real PowerPoint file that opens cleanly in Google Slides.

Usage:  python3 build_abels_board_pptx.py
Output: 2026-06-05-ams-board-presentation.pptx  (in the same folder)

Source of all figures: 2026-06-05-ams-board-presentation.md (confirmed final draft).
HUDL brand: dark theme, four-colour stripe, Barlow Condensed look via Ardc fallbacks.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand palette ─────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x0C, 0x0C, 0x0C)
CARD   = RGBColor(0x14, 0x14, 0x16)
ROW1   = RGBColor(0x0F, 0x0F, 0x12)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFB, 0xBF, 0x46)
BLUE   = RGBColor(0x45, 0xAE, 0xFF)
CYAN   = RGBColor(0x1E, 0xDF, 0xD4)
RED    = RGBColor(0xFF, 0x51, 0x5E)
GREEN  = RGBColor(0x22, 0xC5, 0x5E)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
MUTED  = RGBColor(0x9A, 0x9A, 0x9A)
TBLHDR = RGBColor(0x0A, 0x1A, 0x2A)
FAINT  = RGBColor(0x55, 0x55, 0x55)

FONT_DISPLAY = "Arial Narrow"   # closest universal fallback to Barlow Condensed
FONT_BODY    = "Arial"

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BLACK
    return s

def stripe(s):
    """Four-colour bar across the very top."""
    bw = SW / 4
    for i, col in enumerate([YELLOW, BLUE, RED, CYAN]):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(i*bw)), 0, Emu(int(bw)), Pt(6))
        r.fill.solid(); r.fill.fore_color.rgb = col
        r.line.fill.background()

def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        r.fill.background()
    else:
        r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(line_w)
    r.shadow.inherit = False
    return r

def text(s, txt, l, t, w, h, size=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         font=FONT_BODY, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    lines = txt.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = font
    return tb

def add_table(s, l, t, w, h, headers, rows, col_w=None,
              header_size=11, cell_size=10, cell_colors=None, cell_bold=None,
              row_fills=None):
    nrows, ncols = len(rows) + 1, len(headers)
    gf = s.shapes.add_table(nrows, ncols, l, t, w, h)
    tbl = gf.table
    # remove banding style — we colour manually
    tbl.first_row = False
    tbl.horz_banding = False
    if col_w:
        total = sum(col_w)
        for c, cw in enumerate(col_w):
            tbl.columns[c].width = Emu(int(w * cw / total))
    # header
    for c, htxt in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = TBLHDR
        cell.margin_left = Pt(6); cell.margin_right = Pt(6)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = htxt
        run.font.size = Pt(header_size); run.font.bold = True
        run.font.color.rgb = BLUE; run.font.name = FONT_BODY
    # body
    for r, row in enumerate(rows):
        fill = ROW1 if r % 2 == 0 else BLACK
        if row_fills and r in row_fills:
            fill = row_fills[r]
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = str(val)
            run.font.size = Pt(cell_size); run.font.name = FONT_BODY
            color = WHITE if c == 0 else MUTED
            if cell_colors and r < len(cell_colors) and cell_colors[r] and cell_colors[r][c]:
                color = cell_colors[r][c]
            run.font.color.rgb = color
            b = False
            if cell_bold and r < len(cell_bold) and cell_bold[r] and cell_bold[r][c]:
                b = True
            elif c == 0:
                b = True
            run.font.bold = b
    return tbl

def footer(s, txt="Internal — not for external distribution  ·  HUDL Consultancy  ·  woody@hudl.gg"):
    text(s, txt, Inches(0.4), SH - Inches(0.35), SW - Inches(0.8), Inches(0.25),
         size=8, color=FAINT, align=PP_ALIGN.CENTER)


# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
s = slide()
stripe(s)
rect(s, 0, 0, Pt(6), SH, fill=YELLOW)  # left accent
text(s, "ABELS MOVING SERVICES", Inches(1), Inches(1.7), Inches(11.3), Inches(0.9),
     size=44, color=YELLOW, bold=True, align=PP_ALIGN.CENTER, font=FONT_DISPLAY)
text(s, "Board Marketing Performance Pack", Inches(1), Inches(2.75), Inches(11.3), Inches(0.7),
     size=26, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_DISPLAY)
text(s, "May 2026   |   Jan–May Trend   |   YoY Comparison", Inches(1), Inches(3.5),
     Inches(11.3), Inches(0.5), size=15, color=CYAN, align=PP_ALIGN.CENTER)
rect(s, Inches(4.6), Inches(4.25), Inches(4.13), Pt(2), fill=MUTED)
text(s, "Prepared by HUDL Consultancy — Woody, Fractional Marketing Director\n"
        "Prepared for Neil Pertoldi (AGM Group GM) & AGM Board\n5 June 2026",
     Inches(1), Inches(4.5), Inches(11.3), Inches(1.0), size=12, color=MUTED,
     align=PP_ALIGN.CENTER)
text(s, "CONFIDENTIAL — INTERNAL ONLY", Inches(1), SH - Inches(0.5), Inches(11.3),
     Inches(0.3), size=9, color=FAINT, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — Executive Summary ───────────────────────────────────────────────
s = slide()
stripe(s)
text(s, "EXECUTIVE SUMMARY", Inches(0.4), Inches(0.3), Inches(8), Inches(0.5),
     size=24, color=YELLOW, bold=True, font=FONT_DISPLAY)
text(s, "May 2026", SW - Inches(2.5), Inches(0.42), Inches(2.1), Inches(0.4),
     size=14, color=CYAN, bold=True, align=PP_ALIGN.RIGHT)

# headline callout
rect(s, Inches(0.4), Inches(1.0), SW - Inches(0.8), Inches(0.75), fill=ROW1, line=CYAN, line_w=1.0)
text(s, "Paid media is working. One booked job at 38.5% GM covers ~32–42 leads at the current CPL. "
        "Visitor decline is a deliberate trade of volume for intent — not underperformance.",
     Inches(0.6), Inches(1.08), SW - Inches(1.2), Inches(0.6), size=12, color=WHITE,
     anchor=MSO_ANCHOR.MIDDLE)

rag = {"Green": GREEN, "Amber": AMBER, "Red": RED}
rows = [
    ["Website Visitors",  "2,478 active users",               "Stable vs Apr; spring softening",                  "−36.8% YoY — PMAX removed (strategic)",       "Amber"],
    ["Conversion Volume", "37 lead forms (May)",              "124 Group 2 forms Jan–May (≈24.8/mo)",             "2025: ≈49/mo — different tracking definition","Amber"],
    ["Conversion Quality","38.5% GM · £5,508 job · 21 days",  "GM up 31.3%→38.5%; speed doubled 41→21d",          "Margin and speed both improving YoY",         "Green"],
    ["Cost per Lead",     "£160 blended (May)",               "Holding £160–165 since Apr (Mar spike £228)",      "Genuine 2025 CPL ≈£46 — not like-for-like",   "Green"],
    ["Spend vs Budget",   "£5,929 of ~£6,200 (96%)",          "Near-capacity since Apr (99–96%)",                  "£200/day confirmed by Neil Pertoldi",         "Green"],
]
cell_colors = [[WHITE, MUTED, MUTED, MUTED, rag[r[4]]] for r in rows]
cell_bold   = [[True, False, False, False, True] for _ in rows]
add_table(s, Inches(0.4), Inches(1.95), SW - Inches(0.8), Inches(4.6),
          ["KPI", "May 2026", "Trend (Jan–May)", "YoY context", "Status"],
          rows, col_w=[2.2, 2.6, 3.4, 3.6, 1.3],
          header_size=11, cell_size=10.5, cell_colors=cell_colors, cell_bold=cell_bold)
footer(s)


# ── SLIDES 3–7 — KPI What/Why/How ─────────────────────────────────────────────
kpis = [
    dict(num="01", title="Website Visitors", color=BLUE,
         what="2,478 active users · 2,847 sessions in May\n\n13,329 active users Jan–May\n\n−36.8% YoY (May 2025: 3,918)",
         why="Decline is strategic, not a failure. In 2025, Performance Max generated 11.4 million impressions — almost entirely low-intent display traffic at sub-1% CTR.\n\nWe removed PMAX and switched to five targeted Search-only campaigns. Volume down, intent up.",
         how="9.6% Search conversion rate proves the higher-quality audience.\n\nLevers to recover volume:\n• Unlock Europe's 44% lost impression share\n• Unlock London's 63% lost impression share\n• Structured SEO activity\n\nBoth IS gaps addressed in Q3 (7 July)."),
    dict(num="02", title="Conversion Volume", color=CYAN,
         what="37 lead forms in May\n\n124 Group 2 form submissions Jan–May (≈24.8/mo)\n\nNote: Group 2 = genuine form enquiries only. 2026 counts only this.",
         why="2025 counted soft actions (page views, phone clicks) as conversions — the misleading £21 CPA and ≈177/mo.\n\nGenuine 2025 rate ≈49/mo vs 2026 ≈25/mo. Gap = −37% traffic + lower budget + Domestic inefficiency.",
         how="Europe drove 19 of May's 37 conversions (51%) — clearest signal of where budget belongs.\n\nEurope and London are budget-limited — they stop serving mid-day at their cap. 124 is a floor, not a ceiling. Q3 reallocation targets this."),
    dict(num="03", title="Conversion Quality", color=GREEN,
         what="53 jobs won · £291,919 revenue (Jan–May)\n\n38.5% gross margin (up from 31.3%)\n\n£5,508 average won job\n\n21 days to book (down from 41)\n\nSource: AGM Power BI · all channels.",
         why="The commercial proof: enquiries convert into higher-margin jobs, faster. Gross margin up 31.3%→38.5%. Deals close in half the time (41→21 days). Job value steady ~£5,500 — no trading down.\n\nWinning better jobs, not just more.",
         how="£130–170 CPL against £5,508 job at 38.5% GM (~£2,120 gross profit) = each job covers 12–16 leads at cost.\n\nRouting budget to Europe and London — lowest CPL, highest win rate — lifts monthly revenue without touching margin."),
    dict(num="04", title="Cost per Lead (CPL)", color=YELLOW,
         what="£160 blended (May)\n£165 blended Jan–May\n\nBy campaign (Jan–May):\nEurope £130 · London £151\nInternational £174\nDomestic £299 · Brand £24",
         why="£165 blended is commercially sound (see KPI 3).\n\nDomestic is the outlier: £299/enquiry, the only campaign returning less value than cost (0.74×). Root cause: irrelevant queries — man-and-van, storage, recruitment.",
         how="Domestic Clean-Up Pack (live 23 June):\n• 90 negative keywords\n• Core terms broad → phrase/exact\n• tCPA recalibrated\n\nTarget: £200–230 from £299 over 6–8 weeks. Q3 reallocation channels freed budget into Europe and London."),
    dict(num="05", title="Spend vs Budget", color=CYAN,
         what="May Search: £5,929 (96% of ~£6,200 ceiling)\n\nJan–May Search total: £22,364\n\nTrend: Jan £3,751 · Feb £1,339 · Mar £5,242 · Apr £6,103 · May £5,929\n\nCeiling: £200/day (Neil Pertoldi).",
         why="Account runs at effective capacity Apr–May (99–96%). Goal is efficiency, not raw spend — put each pound where it generates the most profitable enquiries.\n\n£200/day total held in Q3 — reallocation changes distribution, not the total.",
         how="Q3 reallocation (go-live 7 July):\nEurope £97 → £115–130/day\nLondon £28 → £45–55/day\nDomestic £40 → £20/day\nIntl + Brand: unchanged\n\nPhased: clean-up 23 Jun → budgets 7 Jul → tCPA 14 Jul. Smart Bidding needs 2–4 wks to settle."),
]

for k in kpis:
    s = slide()
    stripe(s)
    # badge
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.3),
                               Inches(1.1), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = k["color"]
    badge.line.fill.background(); badge.shadow.inherit = False
    btf = badge.text_frame; btf.word_wrap = False
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = "KPI " + k["num"]
    br.font.size = Pt(13); br.font.bold = True; br.font.color.rgb = BLACK; br.font.name = FONT_BODY
    text(s, k["title"].upper(), Inches(1.7), Inches(0.3), Inches(10), Inches(0.5),
         size=24, color=WHITE, bold=True, font=FONT_DISPLAY)

    colw = (SW - Inches(1.0)) / 3
    xs = [Inches(0.4), Inches(0.4) + colw + Inches(0.1), Inches(0.4) + 2*colw + Inches(0.2)]
    labels = [("WHAT", k["what"]), ("WHY", k["why"]), ("HOW", k["how"])]
    for (lab, body), x in zip(labels, xs):
        text(s, lab, x, Inches(1.05), colw, Inches(0.35), size=12, color=k["color"], bold=True)
        rect(s, x, Inches(1.45), colw, Inches(5.4), fill=ROW1)
        text(s, body, x + Inches(0.15), Inches(1.6), colw - Inches(0.3), Inches(5.1),
             size=11.5, color=MUTED)
    footer(s)


# ── SLIDE 8 — Campaign Activity ───────────────────────────────────────────────
s = slide()
stripe(s)
text(s, "CAMPAIGN ACTIVITY — MAY / JUNE 2026", Inches(0.4), Inches(0.3), Inches(12), Inches(0.5),
     size=22, color=YELLOW, bold=True, font=FONT_DISPLAY)
status_col = {"Completed": GREEN, "Prepared": AMBER, "Recommended": BLUE}
rows = [
    ["Account audit & diagnosis",       "Domestic over-funded; Europe/London under-funded & throttled", "Completed",   "Pippa",         "Done"],
    ["Domestic clean-up pack",          "90 negatives + match-type tightening · CPA £299→£200–230",      "Prepared",    "Pippa",         "23 Jun"],
    ["Budget reallocation",             "Europe ↑ London ↑ Domestic ↓ · total held £200/day",            "Prepared",    "Pippa",         "7 Jul"],
    ["Q3 brief + RSA ad copy",          "15 headlines + 4 descriptions/campaign + 2 landing pages",       "Completed",   "Hunter/Otis",   "Sign-off"],
    ["GA4 analytics review",            "Dashboard updated; quote-path UX opportunity flagged",           "Completed",   "Dash",          "Done"],
    ["Conversion tracking reconciled",  "Proved 2025 counted soft actions; 2026 form-submits only",       "Completed",   "Pippa",         "Done"],
    ["Offline conversion import",       "Re-enable Moveware→Ads (88 leads/£44k in 2025; 0 in 2026)",     "Recommended", "Woody+JDR",     "Q3"],
]
cell_colors = [[WHITE, MUTED, status_col[r[2]], MUTED, MUTED] for r in rows]
cell_bold   = [[True, False, True, False, False] for _ in rows]
add_table(s, Inches(0.4), Inches(1.0), SW - Inches(0.8), Inches(5.6),
          ["Initiative", "What", "Status", "Owner", "Live"], rows,
          col_w=[3.0, 5.4, 1.8, 1.5, 1.2],
          header_size=11, cell_size=10, cell_colors=cell_colors, cell_bold=cell_bold)
footer(s)


# ── SLIDE 9 — Monthly Spend ───────────────────────────────────────────────────
s = slide()
stripe(s)
text(s, "FINANCIAL SUMMARY — MONTHLY SPEND JAN–MAY 2026", Inches(0.4), Inches(0.3),
     Inches(12.5), Inches(0.5), size=20, color=YELLOW, bold=True, font=FONT_DISPLAY)
rows = [
    ["January 2026",  "£3,751",  "29",  "£130",         "~£6,200",            "61%"],
    ["February 2026", "£1,339",  "9",   "£149",         "~£5,600",            "24%"],
    ["March 2026",    "£5,242",  "23",  "£228",         "~£6,200",            "85%"],
    ["April 2026",    "£6,103",  "38",  "£161",         "~£6,000",            "102%"],
    ["May 2026",      "£5,929",  "37",  "£160",         "~£6,200",            "96%"],
    ["Jan–May total", "£22,364", "136", "£165 blended", "£200/day confirmed", "—"],
]
cyan_row   = [CYAN]*6
yellow_row = [YELLOW]*6
cell_colors = [[WHITE, MUTED, MUTED, MUTED, MUTED, MUTED] for _ in rows]
cell_colors[4] = cyan_row
cell_colors[5] = yellow_row
cell_bold = [[False]*6 for _ in rows]
cell_bold[4] = [True]*6
cell_bold[5] = [True]*6
add_table(s, Inches(0.4), Inches(1.1), SW - Inches(0.8), Inches(4.4),
          ["Month", "Search Spend", "Conversions", "CPA", "Budget Ceiling", "Utilisation"], rows,
          col_w=[2.4, 2.0, 2.0, 2.4, 2.6, 1.6],
          header_size=12, cell_size=12, cell_colors=cell_colors, cell_bold=cell_bold,
          row_fills={4: RGBColor(0x0A,0x1A,0x25), 5: RGBColor(0x1A,0x15,0x00)})
text(s, "* April 102% is within daily budget pacing tolerance. February reflects early-year campaign calibration.\n"
        "Total Search (all campaigns) £22,364; JDR-managed campaigns £16,432. Ceiling £200/day confirmed by Neil Pertoldi.",
     Inches(0.4), Inches(5.7), SW - Inches(0.8), Inches(0.8), size=10, color=MUTED)
footer(s)


# ── SLIDE 10 — CPA by Campaign ────────────────────────────────────────────────
s = slide()
stripe(s)
text(s, "CPA BY CAMPAIGN — JAN–MAY 2026 (JDR Search)", Inches(0.4), Inches(0.3),
     Inches(12.5), Inches(0.5), size=20, color=YELLOW, bold=True, font=FONT_DISPLAY)
rows = [
    ["Europe (Outbound)", "£6,119",  "46.89", "£130", "2.57×", "44.4%", "Best in account"],
    ["London",            "£2,872",  "19.00", "£151", "1.83×", "62.9%", "Strong — throttled"],
    ["International",      "£1,391",  "8.00",  "£174", "2.03×", "36.5%", "Efficient — small"],
    ["Brand",             "£72",     "3.00",  "£24",  "5.63×", "2.1%",  "Exceptional"],
    ["Domestic",          "£5,979",  "20.00", "£299", "0.74×", "22.4%", "Requires attention"],
    ["Blended (Search)",  "£16,432", "96.89", "£170", "1.74×", "—",     "Within target"],
]
camp_col = {"Europe (Outbound)": GREEN, "London": CYAN, "International": BLUE,
            "Brand": MUTED, "Domestic": RED, "Blended (Search)": YELLOW}
cell_colors = [[camp_col[r[0]], MUTED, MUTED, MUTED, MUTED, MUTED, MUTED] for r in rows]
cell_colors[4] = [RED]*7
cell_colors[5] = [YELLOW]*7
cell_bold = [[True, False, False, False, False, False, False] for _ in rows]
add_table(s, Inches(0.4), Inches(1.1), SW - Inches(0.8), Inches(3.9),
          ["Campaign", "Spend", "Conv.", "CPA", "Value/Cost", "Lost IS (budget)", "Note"], rows,
          col_w=[2.6, 1.5, 1.3, 1.3, 1.7, 2.2, 2.4],
          header_size=11, cell_size=10.5, cell_colors=cell_colors, cell_bold=cell_bold,
          row_fills={4: RGBColor(0x1A,0x08,0x08), 5: RGBColor(0x1A,0x15,0x00)})
text(s, "MAY 2026 BREAKDOWN", Inches(0.4), Inches(5.25), Inches(4), Inches(0.3),
     size=11, color=MUTED, bold=True)
text(s, "Europe £2,431 · 19 · £128      London £829 · 6 · £138      Domestic £1,196 · 6 · £199      "
        "International £912 · 5 · £182      Brand £49 · 1 · £49",
     Inches(0.4), Inches(5.6), SW - Inches(0.8), Inches(0.6), size=10.5, color=MUTED)
footer(s)


# ── SLIDE 11 — Q3 Roadmap ─────────────────────────────────────────────────────
s = slide()
stripe(s)
text(s, "Q3 ROADMAP", Inches(0.4), Inches(0.3), Inches(5), Inches(0.5),
     size=24, color=YELLOW, bold=True, font=FONT_DISPLAY)
text(s, "Total daily budget held at £200/day throughout", SW - Inches(5.5), Inches(0.45),
     Inches(5.1), Inches(0.4), size=12, color=MUTED, align=PP_ALIGN.RIGHT)

milestones = [
    ("23 Jun", "Domestic clean-up live", "Negatives + match-type tightening", CYAN),
    ("7 Jul",  "Budget reallocation",    "Europe ↑ · London ↑ · Domestic ↓",  YELLOW),
    ("14 Jul", "Domestic tCPA set",      "After clean-up & budgets settle",   BLUE),
]
mw = (SW - Inches(1.0)) / 3
for i, (date, lab, note, col) in enumerate(milestones):
    x = Inches(0.4) + i * (mw + Inches(0.1))
    rect(s, x, Inches(1.15), mw, Inches(1.25), fill=CARD, line=col, line_w=1.5)
    text(s, date, x + Inches(0.15), Inches(1.28), mw - Inches(0.3), Inches(0.4),
         size=16, color=col, bold=True, font=FONT_DISPLAY)
    text(s, lab, x + Inches(0.15), Inches(1.72), mw - Inches(0.3), Inches(0.4),
         size=13, color=WHITE, bold=True)
    text(s, note, x + Inches(0.15), Inches(2.05), mw - Inches(0.3), Inches(0.4),
         size=11, color=MUTED)
text(s, "Smart Bidding needs 2–4 weeks to recalibrate after each change — changes are staggered deliberately.",
     Inches(0.4), Inches(2.55), SW - Inches(0.8), Inches(0.3), size=10, color=FAINT)

# targets (left) + outstanding action (right)
text(s, "Q3 TARGETS", Inches(0.4), Inches(3.05), Inches(4), Inches(0.3),
     size=11, color=MUTED, bold=True)
trows = [
    ["Domestic CPA",          "£299",  "£200–230"],
    ["Blended CPA",           "£165",  "≤ £170 (hold)"],
    ["Europe lost IS budget", "44.4%", "< 20%"],
    ["London lost IS budget", "62.9%", "< 25%"],
]
tcc = [[WHITE, RED, GREEN], [WHITE, MUTED, CYAN], [WHITE, MUTED, CYAN], [WHITE, MUTED, CYAN]]
add_table(s, Inches(0.4), Inches(3.4), Inches(6.0), Inches(2.7),
          ["Metric", "Current", "Q3 Target"], trows, col_w=[2.6, 1.4, 2.0],
          header_size=11, cell_size=11, cell_colors=tcc)

ax = Inches(6.7)
text(s, "OUTSTANDING ACTION", ax, Inches(3.05), Inches(6), Inches(0.3),
     size=11, color=MUTED, bold=True)
rect(s, ax, Inches(3.4), SW - ax - Inches(0.4), Inches(2.7), fill=RGBColor(0x0A,0x10,0x20), line=BLUE, line_w=1.5)
text(s, "Re-enable Moveware → Google Ads offline conversion import", ax + Inches(0.2), Inches(3.55),
     SW - ax - Inches(0.8), Inches(0.6), size=13, color=BLUE, bold=True)
text(s, "2025: 88 qualified leads imported → Smart Bidding optimised toward booked revenue\n"
        "2026: 0 imported → algorithm optimises toward form fills only\n\n"
        "Action: Woody + JDR (upload schedule) + NetDesigner (GCLID capture)\n"
        "Impact: bidding optimises toward revenue, not just enquiries",
     ax + Inches(0.2), Inches(4.15), SW - ax - Inches(0.8), Inches(1.9), size=11, color=MUTED)
footer(s)


# ── SLIDE 12 — Next Steps ─────────────────────────────────────────────────────
s = slide()
stripe(s)
text(s, "NEXT STEPS", Inches(0.4), Inches(0.3), Inches(6), Inches(0.5),
     size=26, color=YELLOW, bold=True, font=FONT_DISPLAY)
steps = [
    ("Sign off Q3 brief & ad copy",        "Europe, London, Domestic RSAs + landing pages ready — needed before 7 July go-live", CYAN,   "Woody → sign off"),
    ("Domestic clean-up: 23 June",         "JDR implements 90 negatives + match-type changes in the Google Ads account",          YELLOW, "JDR Group"),
    ("Budget reallocation: 7 July",        "Europe & London budgets up; Domestic down — total £200/day held",                     YELLOW, "JDR Group"),
    ("Re-enable Moveware offline import",  "Upgrades Smart Bidding from form-fills to revenue optimisation (2025: 88 leads, £44k)",BLUE,   "Woody + JDR + NetDesigner"),
    ("June 2026 board report",             "First full month with reallocation active — expected CPL and lead-volume improvement", MUTED,  "HUDL (Pippa) — July"),
]
y = Inches(1.15)
for lab, note, col, owner in steps:
    rect(s, Inches(0.4), y, SW - Inches(0.8), Inches(0.95), fill=ROW1)
    rect(s, Inches(0.4), y, Pt(4), Inches(0.95), fill=col)
    text(s, lab, Inches(0.65), y + Inches(0.1), Inches(8.0), Inches(0.4),
         size=14, color=WHITE, bold=True)
    text(s, note, Inches(0.65), y + Inches(0.48), Inches(8.5), Inches(0.4),
         size=11, color=MUTED)
    text(s, owner, SW - Inches(3.4), y + Inches(0.27), Inches(3.0), Inches(0.4),
         size=11, color=col, bold=True, align=PP_ALIGN.RIGHT)
    y += Inches(1.12)
footer(s)


# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "2026-06-05-ams-board-presentation.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
